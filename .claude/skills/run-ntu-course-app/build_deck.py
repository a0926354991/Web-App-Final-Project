"""
Build a 6-minute presentation deck for the NTU course recommendation app.

Focus: system features, current progress, todo items.
Avoid: data-scale numbers as the headline.

Outputs to ~/Desktop/Timmy_Report_Output/.

Run:
  .venv/bin/python .claude/skills/run-ntu-course-app/build_deck.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

OUT = Path.home() / "Desktop" / "Timmy_Report_Output"
OUT.mkdir(parents=True, exist_ok=True)

INK = RGBColor(0x0A, 0x0A, 0x0A)
SUB = RGBColor(0x5F, 0x5F, 0x68)
ACCENT = RGBColor(0x00, 0x7A, 0xFF)
DONE = RGBColor(0x21, 0xA1, 0x66)        # green for done
TODO_ACCENT = RGBColor(0xD9, 0x82, 0x00)  # amber for todo
BG = RGBColor(0xFF, 0xFF, 0xFF)

# 16:9
W, H = Emu(12192000), Emu(6858000)


def add_text(slide, x, y, w, h, text, size, color=INK, bold=False, font="Helvetica Neue", align=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = font
    return tb


def add_accent_bar(slide, y=Emu(700000), color=ACCENT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(720000), y, Emu(120000), Emu(80000))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_footer(slide, n, total):
    add_text(slide, Emu(720000), Emu(6300000), Emu(8000000), Emu(400000),
             "台大個性化選課推薦與分析", 10, SUB)
    add_text(slide, Emu(11000000), Emu(6300000), Emu(900000), Emu(400000),
             f"{n} / {total}", 10, SUB)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def section_header(slide, label, color=ACCENT):
    add_accent_bar(slide, color=color)
    add_text(slide, Emu(720000), Emu(820000), Emu(10000000), Emu(700000),
             label, 16, color, bold=True)


# ---------- slides ----------

def make_cover(prs):
    s = blank_slide(prs)
    add_text(s, Emu(720000), Emu(1700000), Emu(10000000), Emu(800000),
             "台大個性化選課推薦", 54, INK, bold=True)
    add_text(s, Emu(720000), Emu(2600000), Emu(10000000), Emu(700000),
             "NTU Personalized Course Recommendation", 28, SUB)
    add_accent_bar(s, y=Emu(3550000))
    add_text(s, Emu(720000), Emu(3700000), Emu(10000000), Emu(400000),
             "把「選課」從靠運氣,變成靠資料。", 20, INK)
    add_text(s, Emu(720000), Emu(5600000), Emu(10000000), Emu(400000),
             "Timmy · Web App 期末專案", 14, SUB)


def make_problem(prs):
    s = blank_slide(prs)
    section_header(s, "01 — 動機")
    add_text(s, Emu(720000), Emu(1300000), Emu(10000000), Emu(900000),
             "選課時最常遇到的三件事。", 32, INK, bold=True)
    items = [
        ("找不到對的課", "課太多、搜尋體驗差,只能靠別人推薦"),
        ("不確定難不難", "PTT 心得散落,沒人有時間翻完"),
        ("不知道適不適合自己", "不同人想要的不一樣,但沒工具會幫你判斷"),
    ]
    y0 = Emu(2700000)
    for i, (title, desc) in enumerate(items):
        y = y0 + i * Emu(820000)
        add_text(s, Emu(720000), y, Emu(420000), Emu(500000),
                 f"0{i+1}", 24, ACCENT, bold=True)
        add_text(s, Emu(1300000), y, Emu(10000000), Emu(500000),
                 title, 22, INK, bold=True)
        add_text(s, Emu(1300000), y + Emu(450000), Emu(10000000), Emu(400000),
                 desc, 14, SUB)
    add_footer(s, 2, 9)


def make_overview(prs):
    s = blank_slide(prs)
    section_header(s, "02 — 系統一頁看懂")
    add_text(s, Emu(720000), Emu(1300000), Emu(10000000), Emu(900000),
             "一個地方,涵蓋選課的整段流程。", 30, INK, bold=True)
    # Four pillars
    pillars = [
        ("瀏覽", "找到想看的課"),
        ("理解", "課綱 + 真實心得一起看"),
        ("規劃", "排課表、列想修、寫筆記"),
        ("推薦", "依你的偏好排序課程"),
    ]
    base_x = Emu(720000)
    col_w = Emu(2700000)
    gap = Emu(180000)
    y = Emu(2900000)
    for i, (title, desc) in enumerate(pillars):
        x = base_x + i * (col_w + gap)
        # rounded card
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, Emu(1700000))
        box.fill.solid(); box.fill.fore_color.rgb = BG
        box.line.color.rgb = INK; box.line.width = Pt(0.75)
        # number
        add_text(s, x + Emu(200000), y + Emu(160000), col_w, Emu(400000),
                 f"0{i+1}", 14, ACCENT, bold=True)
        add_text(s, x + Emu(200000), y + Emu(620000), col_w, Emu(450000),
                 title, 22, INK, bold=True)
        add_text(s, x + Emu(200000), y + Emu(1100000), col_w - Emu(400000), Emu(500000),
                 desc, 13, SUB)
    add_text(s, Emu(720000), Emu(5300000), Emu(11000000), Emu(400000),
             "後端 FastAPI · 前端 Vanilla JS (16 個 ES Modules) · SQLite",
             13, SUB)
    add_footer(s, 3, 9)


def make_features_1(prs):
    s = blank_slide(prs)
    section_header(s, "03 — 功能 ① 找課與看課")
    add_text(s, Emu(720000), Emu(1300000), Emu(10000000), Emu(900000),
             "從搜尋到完整理解一門課。", 30, INK, bold=True)
    feats = [
        ("課程探索", "關鍵字 + 系所 / 學分 / 學期篩選 + 分頁"),
        ("詳情抽屜", "課綱、評量方式、PTT 評價、相關課程一次呈現"),
        ("教師頁面", "點教師名 → 看歷年開課與平均評分"),
        ("課程比較", "勾 2–3 門課並排比較 + 衝堂偵測"),
    ]
    y0 = Emu(2800000)
    for i, (k, v) in enumerate(feats):
        y = y0 + i * Emu(580000)
        # bullet
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(720000), y + Emu(150000), Emu(140000), Emu(140000))
        dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        add_text(s, Emu(1000000), y, Emu(2800000), Emu(500000), k, 18, INK, bold=True)
        add_text(s, Emu(3900000), y + Emu(40000), Emu(8000000), Emu(500000), v, 14, SUB)
    add_footer(s, 4, 9)


def make_features_2(prs):
    s = blank_slide(prs)
    section_header(s, "04 — 功能 ② 個人化規劃")
    add_text(s, Emu(720000), Emu(1300000), Emu(10000000), Emu(900000),
             "讓系統認識你,推薦才會準。", 30, INK, bold=True)
    feats = [
        ("偏好設定", "5 軸能力 + 甜度 / loading 偏好 + 興趣領域"),
        ("修課歷史", "成績 + 個人筆記;已修課自動排除推薦"),
        ("想修清單", "Backend 持久化,登入後跨裝置同步"),
        ("我的課表", "視覺化週課表 + 衝堂提示 + PDF 匯出"),
        ("推薦清單", "Top 20 適合度排序 + 一句話推薦理由"),
    ]
    y0 = Emu(2700000)
    for i, (k, v) in enumerate(feats):
        y = y0 + i * Emu(540000)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(720000), y + Emu(150000), Emu(140000), Emu(140000))
        dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        add_text(s, Emu(1000000), y, Emu(2800000), Emu(500000), k, 18, INK, bold=True)
        add_text(s, Emu(3900000), y + Emu(40000), Emu(8000000), Emu(500000), v, 14, SUB)
    add_footer(s, 5, 9)


def make_algorithm(prs):
    s = blank_slide(prs)
    section_header(s, "05 — 系統差異化:適合度")
    add_text(s, Emu(720000), Emu(1300000), Emu(10000000), Emu(900000),
             "不只給你「別人怎麼想」,還給你「你會不會喜歡」。", 24, INK, bold=True)
    add_text(s, Emu(720000), Emu(2300000), Emu(10000000), Emu(400000),
             "對每門課,從五個面向綜合計算 0–100 分:", 14, SUB)
    parts = [
        ("PTT 推薦度", "結構化評價的推薦指數均值"),
        ("甜度匹配", "你的偏好對上該課的給分風格"),
        ("Loading 匹配", "你的偏好對上該課的負擔程度"),
        ("興趣命中", "13 個興趣 tag 在課程文字裡的 TF-IDF"),
        ("能力匹配", "課程所需能力對上你的長處"),
    ]
    y0 = Emu(2900000)
    for i, (name, desc) in enumerate(parts):
        y = y0 + i * Emu(540000)
        add_text(s, Emu(720000), y, Emu(300000), Emu(500000),
                 "•", 18, ACCENT, bold=True)
        add_text(s, Emu(1000000), y, Emu(3300000), Emu(500000), name, 16, INK, bold=True)
        add_text(s, Emu(4400000), y + Emu(40000), Emu(7400000), Emu(500000), desc, 14, SUB)
    add_text(s, Emu(720000), Emu(5800000), Emu(11000000), Emu(400000),
             "搭配 template 生成的中文推薦理由,讓使用者一眼看懂「為什麼推薦這門」。",
             13, SUB)
    add_footer(s, 6, 9)


def make_progress(prs):
    s = blank_slide(prs)
    section_header(s, "06 — 目前進度", color=DONE)
    add_text(s, Emu(720000), Emu(1300000), Emu(10000000), Emu(900000),
             "核心功能完成,並完成一輪工程化升級。", 26, INK, bold=True)
    sections = [
        ("✓ 已完成", DONE, [
            "完整的探索 / 抽屜 / 教師頁 / 比較 / 課表 / 想修 / 歷史",
            "適合度演算法 5 成份 + 一句話推薦理由",
            "PTT 心得結構化(Claude Batch) + 相關課程 (CF + Content)",
            "深色模式 / RWD / PDF 匯出 / 鍵盤快捷鍵",
        ]),
        ("⚡ 近期重構", ACCENT, [
            "推薦 API 加 TTL 快取(3× 加速)",
            "密碼強度檢查 + 帳號層級 rate limit + production CORS 守門",
            "前端 2008 行單檔拆成 16 個 ES Modules",
            "Playwright E2E 自動化測試 + 截圖驗證",
        ]),
    ]
    y0 = Emu(2500000)
    for i, (title, color, bullets) in enumerate(sections):
        x = Emu(720000) + i * Emu(5800000)
        add_text(s, x, y0, Emu(5500000), Emu(500000), title, 16, color, bold=True)
        for j, b in enumerate(bullets):
            add_text(s, x, y0 + Emu(550000) + j * Emu(550000), Emu(5500000), Emu(500000),
                     f"·  {b}", 13, INK)
    add_footer(s, 7, 9)


def make_todo(prs):
    s = blank_slide(prs)
    section_header(s, "07 — 待做事項", color=TODO_ACCENT)
    add_text(s, Emu(720000), Emu(1300000), Emu(10000000), Emu(900000),
             "從「能用」走向「準確 + 持續活著」。", 28, INK, bold=True)
    items = [
        ("演算法升級", "Sentence Embedding 取代手刻 TF-IDF,提升跨學期推薦準度"),
        ("登入第二方案", "Google OAuth,降低使用者註冊門檻"),
        ("資料持續更新", "APScheduler 每週爬新 PTT 心得,避免推薦準度衰減"),
        ("品質保證", "pytest 單元測試 + Sentry / 結構化日誌"),
        ("更多資料源", "ColleGo、教學意見調查,補 PTT 覆蓋率不足"),
    ]
    y0 = Emu(2600000)
    for i, (k, v) in enumerate(items):
        y = y0 + i * Emu(540000)
        # numbered chip
        chip = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(720000), y, Emu(420000), Emu(420000))
        chip.fill.solid(); chip.fill.fore_color.rgb = TODO_ACCENT
        chip.line.fill.background()
        ctf = chip.text_frame; ctf.text = str(i+1)
        for r in ctf.paragraphs[0].runs:
            r.font.size = Pt(13); r.font.color.rgb = BG; r.font.bold = True
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
        add_text(s, Emu(1300000), y, Emu(2900000), Emu(500000), k, 17, INK, bold=True)
        add_text(s, Emu(4300000), y + Emu(40000), Emu(7500000), Emu(500000), v, 13, SUB)
    add_footer(s, 8, 9)


def make_closing(prs):
    s = blank_slide(prs)
    add_text(s, Emu(720000), Emu(2400000), Emu(10000000), Emu(900000),
             "Demo Time", 60, INK, bold=True)
    add_accent_bar(s, y=Emu(3450000))
    add_text(s, Emu(720000), Emu(3600000), Emu(10000000), Emu(500000),
             "localhost:5500 · 來看看實際運作", 20, SUB)
    add_text(s, Emu(720000), Emu(5600000), Emu(10000000), Emu(400000),
             "Q & A", 18, ACCENT, bold=True)


def build_pptx():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    make_cover(prs)
    make_problem(prs)
    make_overview(prs)
    make_features_1(prs)
    make_features_2(prs)
    make_algorithm(prs)
    make_progress(prs)
    make_todo(prs)
    make_closing(prs)
    path = OUT / "Final_PPT.pptx"
    prs.save(path)
    return path


SCRIPT = """\
# 6-Minute Presentation Script

> 9 張 · 每張約 40 秒 · 總計 6 分鐘
> 主軸:**功能 → 目前進度 → 待做事項**

---

## Slide 1 — 封面 (20 秒)

Hi 大家好,我是 Timmy。
今天要分享的是我的 Web App 期末專案 — **台大個性化選課推薦與分析**。
這個系統想做的事情很簡單:**把「選課」從靠運氣,變成靠資料。**

---

## Slide 2 — 動機 (40 秒)

每個學生在選課時,基本上都會遇到三個問題:

**第一,找不到對的課。** 課太多、搜尋體驗差,
最後只能靠學長姐口耳相傳,但別人推薦的不一定適合你。

**第二,不確定難不難。** PTT 上其實有很多修課心得,
但散落在各個討論串,沒人有時間一篇一篇翻。

**第三,不知道適不適合自己。** 同一門課,
程式底子好的人覺得很涼,沒寫過 code 的人會修到崩潰。
但目前沒有任何工具會幫你判斷「以你的程度,這課適不適合」。

---

## Slide 3 — 系統一頁看懂 (45 秒)

整個系統可以拆成四個動作:

**瀏覽** — 用搜尋跟篩選找想看的課。
**理解** — 課綱 + 真實 PTT 心得一起看,不用切視窗。
**規劃** — 排課表、列想修、寫筆記,全部背景持久化。
**推薦** — 依你填的偏好,排出最適合你的課。

技術上後端 FastAPI、前端 Vanilla JS,
都是相對輕量的選擇,**重點放在使用者體驗而不是技術秀肌肉**。

---

## Slide 4 — 功能 ① 找課與看課 (50 秒)

實際操作分兩塊。第一塊是「找課與看課」:

**課程探索**支援關鍵字 + 系所、學分、學期多重篩選,加上分頁。
點任一門課,**右側抽屜**會把課綱、評量方式、PTT 心得、相關課程一次呈現,
不用跳頁、不用開新分頁。
點教師名,**抽屜會切換到教師頁面**,可以看他歷年所有開過的課跟平均評分。
最後,勾選兩到三門課,**比較頁**會並排呈現所有欄位,
還會自動偵測這幾門有沒有衝堂。

整套體驗刻意做得像 Apple 或 Linear 的 App,**極簡、不囉嗦**。

---

## Slide 5 — 功能 ② 個人化規劃 (50 秒)

第二塊是「個人化規劃」 —— 讓系統認識你。

**偏好設定頁**讓你填五軸能力、甜度跟 loading 偏好、興趣領域。
**修課歷史**可以記成績跟自己的筆記;已修過的課推薦時會自動排除。
**想修清單**是 backend 持久化的,**登入後跨裝置都看得到**。
**我的課表**有視覺化週課表跟衝堂提示,還能一鍵匯出 PDF。
**推薦清單**會列 Top 20,每門課附一句中文推薦理由。

這塊是整個 App 的核心價值 —— 沒有它,就跟一般課程查詢系統一樣。

---

## Slide 6 — 系統差異化:適合度 (50 秒)

跟其他工具的差異就在這:**適合度演算法**。

對每一門課,系統從五個面向綜合計算 0–100 分:
**PTT 推薦度**、**甜度匹配**、**Loading 匹配**、**興趣命中**、**能力匹配**。

PTT 那三個是「別人怎麼想」,後兩個是「你會不會喜歡」。
**ColleGo 給你客觀統計;我的系統給你主觀建議。兩者是互補的。**

每門課除了分數,還會用 template 生出一句中文推薦理由 ——
例如:「PTT 推薦度高、命中你的『金融』興趣,但 loading 比你偏好的重」。

---

## Slide 7 — 目前進度 (55 秒)

進度上,核心功能全部完成,而且這一週做了一輪工程化升級。

**已完成的部分:** 探索、抽屜、教師頁、比較、課表、想修、歷史全部能用;
適合度演算法 5 成份跑通;PTT 心得結構化用 Claude Batch 完成;
相關課程用協同過濾跟內容相似度組成 hybrid 演算法;
還有深色模式、RWD、PDF 匯出、鍵盤快捷鍵這些細節。

**這一週的重構也很大:**
推薦 API 加上 TTL 快取,延遲從 117 毫秒降到 31 毫秒,**三倍加速**;
資安方面加了密碼強度檢查、帳號層級的登入限速、production CORS 守門;
前端 2008 行的單檔拆成 16 個 ES Modules,可維護性大幅提升;
最後寫了一個 Playwright 自動化測試,會跑完整 user flow 並截圖。

---

## Slide 8 — 待做事項 (45 秒)

接下來要做的事,聚焦在「**準確**」跟「**持續活著**」。

第一,**演算法升級** —— 把手刻 TF-IDF 換成 Sentence Embedding,跨學期推薦會更準。
第二,**Google OAuth**,降低使用者註冊摩擦。
第三,**自動排程爬 PTT** —— 現在所有評價都是一次性手動爬,半年後準度會崩。
第四,**單元測試 + 監控告警**,讓品質有保證。
第五,**接更多資料源** —— ColleGo 跟教學意見調查,補 PTT 覆蓋率不足。

這五項做完,系統就具備正式給學弟妹用的條件。

---

## Slide 9 — Demo / Q&A (15 秒)

接下來開實機 demo,順便接受大家的問題。謝謝!
"""

QA = """\
# Possible Q&A

10 個老師或同學最可能問的問題,跟我準備好的回答。

---

**Q1. 為什麼不用 React / Vue,選擇手刻 Vanilla JS?**
A. 三個原因:
這個 App 沒有複雜的 state 管理需求,16 個 ES Module 完全夠用;
不引入框架就不用打包工具,部署成本更低;
也想藉這個機會練「原生 web 技能」,框架會抹掉很多底層細節。

---

**Q2. 你的演算法跟 ColleGo 的差異是什麼?**
A. ColleGo 給「客觀統計」 —— 誰修過、平均成績、分布。
我的系統給「主觀建議」 —— 以你的能力跟偏好,推薦你會喜歡的課。
兩者是互補的:**ColleGo 是字典,這個系統是顧問。**

---

**Q3. 適合度的權重是怎麼決定的?**
A. 老實說是憑直覺加上小規模測試。PTT 評價最直接反映「修過的人怎麼想」,給最高;
甜度、Loading、興趣是「個性化」的核心,各占一塊;
能力匹配的訊號相對弱(很多課關鍵字命中模糊),所以給最低。
**理想上應該用真實使用者點擊數據做 A/B 測試 fit,但這超出期末專案範圍。**

---

**Q4. PTT 心得時效性怎麼處理?半年後不就過期了?**
A. 這正是我「待做事項」清單裡最重要的一項 —— APScheduler 自動排程。
目前是一次性手動爬,大概每隔一兩個月手動跑一次。
真要上線給全校用,**必須改成每週自動爬新貼文**。

---

**Q5. TF-IDF 在中文怎麼處理斷詞?**
A. 我沒做斷詞,直接對 13 個固定 interest tag 做 substring 比對。
英文 tag 用 word boundary(`\\bAI\\b`)避免被 `main`、`again` 誤命中,
中文 tag 直接 substring 找。
tag 集合小且固定,沒必要引入 jieba 這種重量級依賴。

---

**Q6. 用 SQLite 而不是 PostgreSQL,資料量會不會撐不住?**
A. 目前資料量約 30 MB,SQLite 完全夠用。
SQLite 的優點是**零部署成本**(一個檔案)跟**讀效率高**(本專案讀 ≫ 寫)。
如果未來加多人協作課表、即時推薦人數追蹤,才需要升 Postgres。

---

**Q7. 沒有 PTT 心得的課要怎麼推薦?**
A. 對沒有結構化評價的課,推薦時前三項(PTT、甜度、Loading)會用中性 50 分,
這樣分數會偏向「你的偏好」而不是「別人的評價」。
前端也會明確標「⚠ 樣本少」或「無評價」,讓使用者知道訊號弱。

---

**Q8. 密碼用 PBKDF2 不是 bcrypt / argon2,有什麼理由?**
A. PBKDF2-HMAC-SHA256 是 Python 標準函式庫就有,**零依賴**,
20 萬次 iteration 對應目前 OWASP 建議的最低標準。
如果未來要面向真實使用者,會升級到 argon2id,但要多裝一個 C extension。

---

**Q9. 假設要正式上線,你最擔心什麼?**
A. 兩件事:
**PTT 資料時效性** —— 沒有自動排程,推薦準度會慢慢崩;
**惡意爬蟲** —— 有人可能用我的 API 撈整個課程資料庫,
要加 rate limit 跟可能的 captcha。

---

**Q10. 如果有同學想 contribute,你會優先讓他做什麼?**
A. 三個方向都很歡迎:
**寫測試** —— 目前 0 個 pytest,任何 endpoint 都需要;
**改演算法** —— Sentence Embedding 那塊,有 ML 背景的人會很有發揮;
**補資料源** —— 接 ColleGo 或教學意見調查,擴大覆蓋率。
"""


def main():
    pptx_path = build_pptx()
    (OUT / "Presentation_Script.md").write_text(SCRIPT, encoding="utf-8")
    (OUT / "Possible_QA.md").write_text(QA, encoding="utf-8")
    print(f"✅ deck:   {pptx_path}")
    print(f"✅ script: {OUT / 'Presentation_Script.md'}")
    print(f"✅ Q&A:    {OUT / 'Possible_QA.md'}")


if __name__ == "__main__":
    main()
